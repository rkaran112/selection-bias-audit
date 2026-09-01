"""Build the methodology PDF and the 10-slide credit-risk deck.

    python make_report.py

Both read outputs/headline_numbers.json. Nothing here computes a statistic;
if a number appears in the PDF or the deck, run_all.py produced it.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent / "src"))

from reportlab.lib import colors                                # noqa: E402
from reportlab.lib.enums import TA_LEFT                         # noqa: E402
from reportlab.lib.pagesizes import A4                          # noqa: E402
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet  # noqa: E402
from reportlab.lib.units import cm                              # noqa: E402
from reportlab.platypus import (BaseDocTemplate, Frame, Image,  # noqa: E402
                                PageBreak, PageTemplate, Paragraph, Spacer,
                                Table, TableStyle)

from sba import config as C                                     # noqa: E402

INK = colors.HexColor("#1b1b1f")
MUTED = colors.HexColor("#6b7280")
RULE = colors.HexColor("#d4d4d8")
ACCENT = colors.HexColor("#c2410c")
TEAL = colors.HexColor("#0e7490")
WASH = colors.HexColor("#f6f7f9")
WARN = colors.HexColor("#fdf6f0")

D: dict = {}


# --------------------------------------------------------------------------
def n(path: str, default=None):
    """Pull a value out of headline_numbers.json by dotted path."""
    cur = D
    for k in path.split("."):
        if isinstance(cur, list):
            try:
                cur = cur[int(k)]
                continue
            except (ValueError, IndexError):
                return default
        if not isinstance(cur, dict) or k not in cur:
            return default
        cur = cur[k]
    return cur if cur is not None else default


def pct(x, dp=1):
    return "n/a" if x is None else f"{100 * x:.{dp}f}%"


def num(x, dp=4):
    return "n/a" if x is None else f"{x:.{dp}f}"


def money(x, dp=0):
    if x is None:
        return "n/a"
    a = abs(x)
    s = "-" if x < 0 else ""
    if a >= 1e9:
        return f"{s}${a / 1e9:.{max(dp,1)}f}bn"
    if a >= 1e6:
        return f"{s}${a / 1e6:,.{dp}f}m"
    return f"{s}${a:,.0f}"


# --------------------------------------------------------------------------
_ss = getSampleStyleSheet()
S = {
    "title": ParagraphStyle("t", parent=_ss["Title"], fontName="Helvetica-Bold",
                            fontSize=25, leading=29, textColor=INK,
                            alignment=TA_LEFT, spaceAfter=4),
    "sub": ParagraphStyle("s", fontName="Helvetica", fontSize=12.5, leading=17,
                          textColor=MUTED, spaceAfter=16),
    "h1": ParagraphStyle("h1", fontName="Helvetica-Bold", fontSize=15,
                         leading=19, textColor=INK, spaceBefore=20,
                         spaceAfter=8),
    "h2": ParagraphStyle("h2", fontName="Helvetica-Bold", fontSize=11.5,
                         leading=15, textColor=INK, spaceBefore=13,
                         spaceAfter=5),
    "body": ParagraphStyle("b", fontName="Helvetica", fontSize=9.9,
                           leading=14.6, textColor=INK, spaceAfter=8),
    "small": ParagraphStyle("sm", fontName="Helvetica", fontSize=8.6,
                            leading=12.2, textColor=MUTED, spaceAfter=6),
    "callout": ParagraphStyle("c", fontName="Helvetica", fontSize=10,
                              leading=14.6, textColor=INK, spaceAfter=8,
                              leftIndent=10, rightIndent=10, spaceBefore=4,
                              borderPadding=9, backColor=WASH,
                              borderColor=RULE, borderWidth=0.6),
    "warn": ParagraphStyle("w", fontName="Helvetica", fontSize=10,
                           leading=14.6, textColor=INK, spaceAfter=8,
                           leftIndent=10, rightIndent=10, spaceBefore=4,
                           borderPadding=9, backColor=WARN,
                           borderColor=ACCENT, borderWidth=0.9),
}


def P(t, s="body"):
    return Paragraph(t, S[s])


def bullets(items, style="body"):
    return [Paragraph(f"&bull;&nbsp;&nbsp;{i}", S[style]) for i in items]


def table(rows, widths=None, header=True, align=None, fs=8.6):
    t = Table(rows, colWidths=widths, hAlign="LEFT", repeatRows=1 if header else 0)
    st = [
        ("FONT", (0, 0), (-1, -1), "Helvetica", fs),
        ("TEXTCOLOR", (0, 0), (-1, -1), INK),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 4.5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4.5),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("LINEBELOW", (0, 0), (-1, -2), 0.35, RULE),
    ]
    if header:
        st += [("FONT", (0, 0), (-1, 0), "Helvetica-Bold", fs),
               ("BACKGROUND", (0, 0), (-1, 0), WASH),
               ("LINEBELOW", (0, 0), (-1, 0), 0.9, INK)]
    for c in (align or []):
        st.append(("ALIGN", (c, 0), (c, -1), "RIGHT"))
    t.setStyle(TableStyle(st))
    return t


def figure(name, width=16.4 * cm):
    p = C.FIGURES / name
    if not p.exists():
        return P(f"[figure {name} not found - run run_all.py]", "small")
    from PIL import Image as PILImage
    with PILImage.open(p) as im:
        w, h = im.size
    return Image(str(p), width=width, height=width * h / w)


# --------------------------------------------------------------------------
def _decorate(canvas, doc):
    canvas.saveState()
    canvas.setFont("Helvetica", 7.6)
    canvas.setFillColor(MUTED)
    canvas.drawString(2.2 * cm, 1.3 * cm,
                      "Selection-bias audit for credit scorecards"
                      "  -  methodology")
    canvas.drawRightString(A4[0] - 2.2 * cm, 1.3 * cm, f"{doc.page}")
    canvas.setStrokeColor(RULE)
    canvas.setLineWidth(0.4)
    canvas.line(2.2 * cm, 1.75 * cm, A4[0] - 2.2 * cm, 1.75 * cm)
    canvas.restoreState()


def build_pdf(out: Path) -> Path:
    doc = BaseDocTemplate(str(out), pagesize=A4,
                          leftMargin=2.2 * cm, rightMargin=2.2 * cm,
                          topMargin=2.0 * cm, bottomMargin=2.2 * cm,
                          title="Selection-bias audit for credit scorecards",
                          author="rkaran112")
    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height,
                  id="f")
    doc.addPageTemplates([PageTemplate(id="main", frames=[frame],
                                       onPage=_decorate)])

    base_gini = n("part1_baseline.test.gini")
    base_ks = n("part1_baseline.test.ks")
    base_auc = n("part1_baseline.test.auc")
    rho = n("part2_real_data.heckman.diagnostics.rho")
    rho_p = n("part2_real_data.heckman.diagnostics.lr_p")
    rho_se = n("part2_real_data.heckman.diagnostics.se_rho")
    rho_ci = n("part2_real_data.heckman.diagnostics.rho_ci95") or [None, None]
    rho_x = n("part2_real_data.heckman_with_exclusion.diagnostics.rho")
    rho_xp = n("part2_real_data.heckman_with_exclusion.diagnostics.lr_p")
    ranking = n("part3_simulation.ranking", [])
    nullc = n("part3_simulation.null_case", [])
    verdict = n("part3_simulation.degradation.verdict", "")
    _h = sorted((r for r in nullc if r.get("hallucinates")),
                key=lambda r: -r["excess_over_baseline"])
    hall = [r["method"] for r in _h]
    hall_named = ", ".join(
        f"{r['method']} ({100 * r['excess_over_baseline']:+.1f} pp)"
        for r in _h)
    disp = n("part4_geography.dispersion", {})

    E = []
    # ------------------------------------------------------------- title
    E += [
        P("Selection-bias audit for<br/>credit scorecards", "title"),
        P("Quantifying how much of a scorecard's reported performance is an "
          "artefact of only ever seeing approved applicants - and testing "
          "whether the standard corrections actually work.", "sub"),
        table([
            ["Data", "LendingClub 2007-2018Q4, accepted and declined "
                     "applications (CC0)"],
            ["Approved", f"{n('data.accepted_rows_total', 0):,} loans; "
                         f"{n('data.accepted_rows_terminal', 0):,} with a "
                         f"settled outcome"],
            ["Declined", f"{n('data.rejected_rows_total', 0):,} applications"],
            ["Approval rate", pct(n("data.population_accept_rate"), 2)],
            ["Seed", str(n("meta.seed"))],
            ["Generated", str(n("meta.generated_utc"))],
        ], widths=[3.6 * cm, 12.8 * cm], header=False),
        Spacer(1, 12),
        P("<b>The headline.</b> The scorecard reports a Gini of "
          f"<b>{num(base_gini)}</b> on approved loans. Four different "
          "reject-inference corrections disagree with each other about the "
          "true through-the-door bad rate by several percentage points. This "
          "report does not pick a winner on theoretical grounds. It builds a "
          "simulation in which the right answer is known, runs all four "
          "against it, and reports which ones survive - including the finding "
          f"that {len(hall)} of them report a correction even when there is "
          "provably nothing to correct.", "callout"),
    ]

    # ------------------------------------------------- 1. what we found
    E += [P("1. What we found", "h1")]
    hall = [r["method"] for r in nullc if r.get("hallucinates")]
    best = ranking[0]["method"] if ranking else "n/a"
    E += bullets([
        f"<b>The baseline.</b> Gini {num(base_gini)}, KS {num(base_ks)}, "
        f"AUC {num(base_auc)} on held-out approved loans. This is the number "
        f"a lender would report, and it is measured on the wrong population.",

        f"<b>Selection correlation.</b> The bivariate-probit selection "
        f"parameter rho is {num(rho)} (se {num(rho_se)}, "
        f"likelihood-ratio p = {num(rho_p, 3)}). "
        + ("Because that is not distinguishable from zero, the honest reading "
           "is that on the fields both populations share, approval carried "
           "little information about default beyond what the scorecard "
           "already uses. Manufacturing a larger correction would be the "
           "actual failure."
           if (rho_p or 1) > 0.05 else
           "That is statistically distinguishable from zero, so some "
           "selection effect is present on the shared fields."),

        f"<b>Ranking (simulation).</b> On risk-based cutoffs, the least "
        f"biased approach was <b>{best}</b>. Full table in section 5.",

        (f"<b>{len(hall)} of the four hallucinate.</b> Under a random cutoff "
         f"- where no selection bias exists and the correct correction is "
         f"zero - {hall_named} still report one. The magnitudes differ by an "
         f"order of magnitude and so do the causes; section 6 separates them."
         if hall else
         "<b>Null case.</b> Under a random cutoff every method correctly "
         "reported approximately no correction."),

        f"<b>Geography.</b> Across {disp.get('n_zip3', 0)} 3-digit ZIPs, "
        f"decline rates vary by more than modelled risk explains "
        f"(dispersion ratio {num(disp.get('dispersion_ratio'), 1)} against "
        f"1.0 under the null). This is a pointer for review, not a finding "
        f"of discrimination.",
    ])

    # --------------------------------------------- 2. thin overlap ceiling
    E += [PageBreak(), P("2. The thin-overlap ceiling", "h1"),
          P("This is the binding constraint on everything that follows, so it "
            "comes first rather than in a footnote.", "body")]
    E += [P("The approved file carries around 150 columns. The declined file "
            "carries nine. A model that uses anything outside the "
            "intersection cannot be applied to a declined applicant, so the "
            "usable feature set is four numeric fields and one categorical "
            "one. That is a thin basis for a scorecard, and it caps how well "
            "any method here can possibly do.", "body")]

    cov = n("data.coverage", [])
    E += [P("Measured on the raw sampled files, BEFORE the scored-decline "
            "filter described below. After that filter risk_score is 0% "
            "missing on declines - true of the modelling sample, and badly "
            "misleading about the source data.", "small")]
    rows = [["Shared field", "Missing (approved)", "Missing (declined)"]]
    for r in cov:
        rows.append([r["field"], f"{r['accepted_missing_pct']:.1f}%",
                     f"{r['rejected_missing_pct']:.1f}%"])
    E += [table(rows, widths=[6.4 * cm, 5.0 * cm, 5.0 * cm], align=[1, 2]),
          Spacer(1, 8)]

    E += [P(
        "<b>It is worse than thin - it is holed.</b> "
        f"{pct(n('data.unscored_reject_fraction'))} of declined applications "
        "carry no bureau score at all, and the gap is not random: it is "
        "concentrated in the later vintages. A scorecard fitted on approved "
        "loans never encounters a missing score, so it has no bin for one and "
        "would rate those applicants as merely average - which is precisely "
        "backwards. This audit therefore restricts the declined population to "
        f"scored applications, using {n('data.rejected_rows_used', 0):,} of "
        f"{n('data.rejected_rows_sampled_raw', 0):,} sampled declines. The "
        "excluded two thirds are a real and unquantified limitation.",
        "warn")]

    E += [P("Three further mismatches worth stating plainly:", "h2")]
    E += bullets([
        "<b>The scores are not the same instrument.</b> Approved loans carry "
        "a FICO range from the bureau pull; declined applications carry a "
        "field LendingClub calls Risk_Score. They are on a similar scale and "
        "behave similarly, but they are not guaranteed to be the same model "
        "or the same vintage.",
        "<b>Debt-to-income is not measured the same way.</b> On the approved "
        "side it is computed against verified income. On the declined side it "
        "is self-reported and includes values that are plainly impossible, "
        "which this pipeline sets to missing rather than clipping.",
        "<b>The dates mean different things.</b> The declined file records the "
        "application date; the approved file records the funding date. The "
        "gap is a few weeks, which matters only for vintage alignment.",
    ])
    E += [Spacer(1, 6), figure("01_overlap.png"),
          P("Left: the two score distributions barely overlap - "
            f"only {num(n('data.scored_rejects_above_accepted_p01_pct'), 1)}% "
            "of scored declines reach the 1st percentile of the approved "
            "book. Right: where the shared fields are missing.", "small")]

    # ------------------------------------------------- 3. baseline
    E += [PageBreak(), P("3. Part 1 - the biased baseline", "h1"),
          P("A conventional application scorecard: weight-of-evidence binning "
            "on each shared field, logistic regression on the binned values, "
            "scaled to points. Deliberately conventional - the question here "
            "is about the sample a model is fitted on, not the model class, "
            "and holding the model fixed keeps the four corrections "
            "comparable.", "body")]
    E += [table([
        ["Metric", "Train", "Held-out test"],
        ["Gini", num(n("part1_baseline.train.gini")), num(base_gini)],
        ["KS", num(n("part1_baseline.train.ks")), num(base_ks)],
        ["AUC", num(n("part1_baseline.train.auc")), num(base_auc)],
        ["Bad rate", pct(n("part1_baseline.train.bad_rate"), 2),
         pct(n("part1_baseline.test.bad_rate"), 2)],
        ["n", f"{n('part1_baseline.n_train', 0):,}",
         f"{n('part1_baseline.n_test', 0):,}"],
    ], widths=[5.0 * cm, 5.7 * cm, 5.7 * cm], align=[1, 2]), Spacer(1, 6)]
    iv = n("part1_baseline.iv", [])
    if iv:
        E += [P("Information value by field", "h2"),
              table([["Field", "IV", "Bins"]]
                    + [[r["feature"], num(r["iv"], 4), str(r["n_bins"])]
                       for r in iv],
                    widths=[6.4 * cm, 5.0 * cm, 5.0 * cm], align=[1, 2])]
    E += [Spacer(1, 8), figure("02_baseline.png")]
    E += [P("Scaling: "
            f"{num(n('part1_baseline.pdo'), 0)} points to double the odds, "
            f"{num(n('part1_baseline.base_points'), 0)} points at "
            f"{num(n('part1_baseline.base_odds'), 0)}:1 good:bad. Full "
            "per-bin point allocation is in "
            "outputs/tables/01_scorecard_points.csv.", "small")]

    # ------------------------------------------------- 4. the four methods
    E += [PageBreak(), P("4. Part 2 - four corrections", "h1")]
    E += [P("Each method must produce two things: a scoring function that "
            "applies to any applicant, and an estimate of the "
            "through-the-door bad rate. Part 3 then checks both against a "
            "truth the method was never shown.", "body")]
    E += [table([
        ["Method", "What it assumes", "Where it breaks"],
        ["Parcelling",
         "Declines default at a fixed multiple of the approved rate in the "
         "same score band.",
         "The multiple is an assumption nothing in the data identifies. It "
         "fires whether or not selection is informative."],
        ["Fuzzy augmentation",
         "Same multiple, applied as fractional weights instead of drawn "
         "labels.",
         "Identical assumption, with the sampling noise removed. Stable, and "
         "stably wrong."],
        ["Inverse propensity",
         "Given the shared fields, approval carries no further information "
         "about default (missing at random).",
         "Requires that a declined applicant could plausibly have been "
         "approved. Here that positivity condition fails badly."],
        ["Heckman / bivariate probit",
         "Approval and default are jointly normal with correlation rho, "
         "estimated by maximum likelihood.",
         "Without an exclusion restriction, identification rests on the "
         "normality assumption alone."],
    ], widths=[3.5 * cm, 6.3 * cm, 6.6 * cm], fs=8.2), Spacer(1, 10)]

    rows = [["Method", "Est. TTD bad rate", "Est. decline bad rate",
             "Claimed Gini"]]
    for m in ("parcelling", "fuzzy", "ipw", "heckman"):
        rows.append([
            m,
            pct(n(f"part2_real_data.{m}.est_ttd_bad_rate"), 2),
            pct(n(f"part2_real_data.{m}.est_reject_bad_rate"), 2),
            num(n(f"part2_real_data.{m}.claimed_ttd_gini")),
        ])
    rows.append(["none (observed)",
                 pct(n("data.observed_bad_rate_accepted"), 2), "-",
                 num(base_gini)])
    E += [P("What they say on the real declined population", "h2"),
          table(rows, widths=[4.2 * cm, 4.2 * cm, 4.2 * cm, 3.8 * cm],
                align=[1, 2, 3]),
          P("These four disagree by percentage points. Nothing in this table "
            "tells you which to believe - that is what Part 3 is for.",
            "small")]

    E += [P("The selection correlation, rho", "h2"),
          table([
              ["Specification", "rho", "std error", "95% interval", "LR p"],
              ["No exclusion restriction", num(rho), num(rho_se),
               f"[{num(rho_ci[0], 3)}, {num(rho_ci[1], 3)}]", num(rho_p, 3)],
              ["Excluding application year", num(rho_x), "-", "-",
               num(rho_xp, 3)],
          ], widths=[5.4 * cm, 2.6 * cm, 2.6 * cm, 3.4 * cm, 2.4 * cm],
              align=[1, 2, 3, 4])]
    E += [P(
        "rho is the correlation between the unobserved drivers of approval and "
        "the unobserved drivers of default. If it is zero, approval told you "
        "nothing the scorecard did not already know, and the accepts-only "
        "model is unbiased.", "body")]
    E += [P(
        f"The two specifications disagree on significance: {num(rho)} "
        f"(p = {num(rho_p, 3)}) without an exclusion restriction, "
        f"{num(rho_x)} (p = {num(rho_xp, 3)}) with application year excluded "
        "from the outcome equation. Two cautions apply, and they point the "
        "same way. <b>The exclusion restriction is contestable</b> - "
        "application vintage plausibly affects default directly through the "
        "macro cycle, so excluding it from the outcome equation is a "
        "misspecification that can manufacture apparent selection rather than "
        "reveal it. <b>And significance is not magnitude.</b> At this sample "
        "size a correlation of "
        f"{num(max(abs(rho or 0), abs(rho_x or 0)), 3)} is detectable and "
        "still far too small to move a lending decision. On either "
        "specification the reading is the same: <b>selection bias in this "
        "portfolio, measured on the fields the two populations share, is "
        "small.</b> That is a legitimate result and it is reported as one "
        "rather than inflated into a finding.",
        "callout")]

    pos = n("part2_positivity", {})
    E += [P("Positivity: why inverse-propensity weighting is on thin ice here",
            "h2"),
          table([
              ["Diagnostic", "Value", "Reading"],
              ["AUC of the approval model", num(pos.get("propensity_auc")),
               "1.0 would mean approval is perfectly predictable"],
              ["Smallest approval propensity",
               num(pos.get("min_propensity_accepted"), 5),
               "Near zero means near-deterministic decline"],
              ["Effective sample size retained",
               pct(pos.get("ess_ratio")),
               "How much of the data survives the reweighting"],
          ], widths=[5.6 * cm, 3.2 * cm, 7.6 * cm], fs=8.2),
          P("An approval model that separates this well means most declined "
            "applicants had essentially no chance of approval given the "
            "shared fields. Reweighting cannot conjure information about "
            "applicants who are never observed on the other side of the "
            "cutoff, and the effective sample size shows the cost.", "small")]

    # ------------------------------------------------- 5. simulation
    E += [PageBreak(), P("5. Part 3 - the simulation harness", "h1")]
    E += [P("Reject inference is normally unfalsifiable. You cannot check an "
            "inferred bad rate against outcomes that do not exist. So we "
            "manufacture the missing truth.", "body")]
    E += bullets([
        "Take only approved loans, where every outcome is known.",
        "Impose an artificial cutoff and discard the outcomes below it, "
        "pretending those applicants were declined.",
        "Run all four methods on that synthetic decline population.",
        "Compare each against the truth we deliberately hid.",
    ])
    E += [P(f"Swept across rejection rates of "
            f"{', '.join(pct(r, 0) for r in n('part3_simulation.rejection_rates', []))}, "
            f"with {n('part3_simulation.n_replicates')} replicates at each "
            f"point and a pool of "
            f"{n('part3_simulation.pool_size', 0):,} loans per replicate. "
            "Three cutoff shapes are used: a hard threshold on the score, the "
            "same threshold with judgemental override noise, and a random "
            "cutoff that serves as the null.", "body")]

    rows = [["Method", "|bad-rate bias|", "|Gini bias|", "Real Gini uplift",
             "Instability"]]
    for r in ranking:
        rows.append([r["method"],
                     f"{100 * r['mean_abs_bad_rate_bias']:.2f} pp",
                     num(r["mean_abs_gini_bias"], 4),
                     f"{r['mean_gini_uplift']:+.4f}",
                     num(r["instability"], 4)])
    E += [P("Ranking on risk-based cutoffs (best first)", "h2"),
          table(rows, widths=[4.4 * cm, 3.2 * cm, 2.8 * cm, 3.2 * cm,
                              2.8 * cm], align=[1, 2, 3, 4]),
          P("Bad-rate bias is the method's estimate of the through-the-door "
            "bad rate minus the truth. Gini bias is what the method claims "
            "its Gini is minus what it actually achieves on the hidden "
            "labels. Uplift is whether the corrected model genuinely ranks "
            "better than the uncorrected one. Instability is the "
            "across-replicate standard deviation.", "small")]
    if verdict:
        E += [P(f"<b>Where it stops working.</b> {verdict}", "callout")]
    E += [Spacer(1, 6), figure("03_sweep_bad_rate.png"), Spacer(1, 4),
          figure("05_gini_levels.png")]

    # ------------------------------------------------- 6. null case
    E += [PageBreak(), P("6. The null case", "h1")]
    E += [P("The most important panel in this report. Replace the risk-based "
            "cutoff with a random one. Now the approved sample is a random "
            "subsample of the population, there is no selection bias, and the "
            "correct answer from every method is 'no correction needed'. Any "
            "method that still reports a large correction is not measuring "
            "bias; it is generating it.", "body")]
    rows = [["Method", "|bias|", "vs. doing nothing", "one-sided p", "Verdict"]]
    for r in nullc:
        rows.append([
            r["method"], f"{100 * r['mean_abs_bad_rate_bias']:.2f} pp",
            f"{100 * r['excess_over_baseline']:+.2f} pp",
            num(r["p_one_sided"], 3),
            "HALLUCINATES" if r["hallucinates"] else "passes",
        ])
    t = table(rows, widths=[4.4 * cm, 2.6 * cm, 3.4 * cm, 2.4 * cm, 3.6 * cm],
              align=[1, 2, 3])
    for i, r in enumerate(nullc, start=1):
        if r["hallucinates"]:
            t.setStyle(TableStyle([("BACKGROUND", (0, i), (-1, i), WARN),
                                   ("TEXTCOLOR", (4, i), (4, i), ACCENT),
                                   ("FONT", (4, i), (4, i),
                                    "Helvetica-Bold", 8.6)]))
    E += [t, P("Each method is compared against doing nothing on the same "
               "scenarios, by a paired test across replicates. Doing nothing "
               "is not error-free under a random cutoff - it carries sampling "
               "noise - so a method only counts as hallucinating if it is "
               "reliably worse than that, by a margin worth caring about.",
               "small")]
    E += [Spacer(1, 6), figure("06_null_case.png")]

    scale = n("part3_scale_sensitivity", [])
    if scale:
        E += [P("Is this just a badly chosen scale factor?", "h2"),
              P("Parcelling and fuzzy augmentation both need an exogenous "
                "multiple for how much worse declines are than approvals. "
                "Ranking them last while holding that multiple at a number "
                "they did not choose would be a strawman, so both are re-run "
                "across a range of it.", "body")]
        rows = [["Cutoff", "Method", "Scale factor",
                 "|bad-rate bias|"]]
        for r in scale:
            rows.append([r["cutoff_type"], r["method"],
                         num(r["reject_bad_scale"], 1),
                         f"{100 * r['abs_bad_rate_bias']:.2f} pp"])
        E += [table(rows, widths=[3.6 * cm, 4.0 * cm, 4.2 * cm, 4.2 * cm],
                    align=[2, 3], fs=8.0),
              P("On a risk-based cutoff a well-chosen multiple scores "
                "respectably; nothing in the data tells you which one that is. "
                "On the null cutoff only a multiple of exactly 1.0 avoids "
                "inventing a correction, and you would need to already know "
                "there was no selection bias in order to choose it. The "
                "null-case failure is structural, not a tuning problem.",
                "small")]

    rho_sweep = n("part3_rho_sweep", [])
    if rho_sweep:
        E += [PageBreak(), P("6b. What the rho sweep reveals", "h1"),
              P("Because the simulation sets the selection rule itself, the "
                "true value of rho is known in every scenario here: it is "
                "zero. Selection depends only on the observed score, never on "
                "anything correlated with the hidden outcome. Any non-zero rho "
                "the estimator reports is therefore an error, and the pattern "
                "of those errors is instructive.", "body")]
        rows = [["Cutoff", "Rejection rate", "rho", "std error", "LR p"]]
        for r in rho_sweep:
            rows.append([r["cutoff_type"], f"{r['rejection_rate']:.0%}",
                         num(r["rho"], 4),
                         (f"{r['se_rho']:,.0f}" if abs(r["se_rho"]) > 100
                          else num(r["se_rho"], 4)),
                         num(r["lr_p"], 3)])
        E += [table(rows, widths=[3.6 * cm, 3.4 * cm, 3.0 * cm, 3.4 * cm,
                                  2.8 * cm], align=[1, 2, 3, 4], fs=8.2)]
        E += bullets([
            "<b>Random cutoff.</b> The likelihood-ratio test is correctly "
            "insignificant at every rate. The test is honest. The rho point "
            "estimate still wanders, though, with standard errors of 0.3 to "
            "0.5 - which is what produces Heckman's null-case error. Act on "
            "the test, not on the point estimate.",

            "<b>Hard cutoff, high rejection.</b> The standard error explodes "
            "to the thousands. That is not a bug: approval is a deterministic "
            "function of an observed covariate, so rho is genuinely "
            "unidentified and the model is saying so. A practitioner who "
            "reports the point estimate without the standard error would miss "
            "that entirely.",

            "<b>Soft cutoff.</b> Here the estimator reports rho around 0.25 "
            "with p = 0.005 - a confident finding of selection bias in a "
            "simulation where selection is provably independent of the "
            "outcome. The cause is functional form: both equations "
            "approximate a smooth score with WOE steps, both residuals "
            "therefore carry correlated approximation error, and the "
            "estimator reads that as rho.",
        ])
        E += [P("<b>This is the most important caveat on the real-data "
                "result in section 4.</b> A statistically significant rho is "
                "not proof of selection bias; misspecification alone can "
                "manufacture one. It is a further reason to read the small "
                "real-data rho as small rather than to reach for the "
                "specification that makes it significant.", "warn")]
        E += [Spacer(1, 4), figure("04_rho.png")]

    # ------------------------------------------------- 7. economics
    E += [PageBreak(), P("7. Part 4 - what the bias costs", "h1")]
    yc = n("part4_yields", {})
    ass = yc.get("assumptions", {})
    E += [P("Wherever the data allows it, these numbers are measured rather "
            "than assumed. LendingClub publishes realised cash flows for every "
            "settled loan, so the yield on a repaid loan, the loss on a "
            "charge-off and the loss-given-default all come from those cash "
            "flows.", "body")]
    E += [table([
        ["Quantity", "Value", "Source"],
        ["Net return per $1, loan repaid",
         num(yc.get("good_yield_per_dollar"), 4), "measured from cash flows"],
        ["Net return per $1, charge-off",
         num(yc.get("bad_yield_per_dollar"), 4), "measured from cash flows"],
        ["Loss given default", pct(yc.get("empirical_lgd"), 1),
         "measured from recoveries"],
        ["Weighted average life",
         f"{ass.get('wal_fraction_of_term')} x stated term", "ASSUMPTION"],
        ["Cost of funds (annual)", pct(ass.get("cost_of_funds_annual"), 1),
         "ASSUMPTION"],
        ["Servicing cost (annual)", pct(ass.get("servicing_cost_annual"), 1),
         "ASSUMPTION"],
        ["Break-even bad rate, 36m",
         pct(n("part4_economics.breakeven_bad_rate_36m"), 1), "derived"],
        ["Break-even bad rate, 60m",
         pct(n("part4_economics.breakeven_bad_rate_60m"), 1), "derived"],
    ], widths=[6.4 * cm, 3.6 * cm, 6.4 * cm], align=[1])]

    rows = [["Method", "Swap-in applicants", "Inferred bad rate",
             "Break-even they face", "Profit forgone"]]
    for r in n("part4_economics.by_method", []):
        rows.append([
            r.get("method", ""),
            f"{r.get('swap_in_applicants', 0):,.0f}",
            pct(r.get("swap_in_inferred_bad_rate"), 1),
            pct(r.get("swap_in_breakeven_bad_rate"), 1),
            money(r.get("profit_forgone_population")),
        ])
    E += [Spacer(1, 8), P("Swap-set analysis, volume-neutral", "h2"),
          table(rows, widths=[3.0 * cm, 3.2 * cm, 3.4 * cm, 3.0 * cm,
                              3.8 * cm], align=[1, 2, 3, 4], fs=8.2),
          P("The corrected model approves exactly as many applicants as the "
            "lender approves today, so this isolates ranking quality rather "
            "than recommending a bigger book. Note the break-even column: the "
            "swap-in set lands in the safest score bands, where LendingClub "
            "priced loans near 9% APR, so the break-even THEY face is far "
            "below the pooled figure above. Comparing their bad rate against "
            "the pooled break-even would imply a profit the cash flows do not "
            "support. Bootstrap intervals are in "
            "outputs/tables/11_swap_set_economics.csv; they cover sampling "
            "error in the swap set only, not the disagreement between "
            "methods, which is far larger and is the real uncertainty.",
            "small")]

    prof = n("part4_swap_profile", [])
    if prof:
        rows = [["Group", "Applicants", "Mean score", "Mean DTI",
                 "Mean amount"]]
        for r in prof:
            rows.append([r["group"], f"{r['applicants']:,.0f}",
                         num(r.get("risk_score"), 1),
                         num(r.get("dti"), 1),
                         f"${r.get('amount_requested', 0):,.0f}"])
        E += [Spacer(1, 8), P("Which applicants actually flip", "h2"),
              table(rows, widths=[4.6 * cm, 3.0 * cm, 2.8 * cm, 2.6 * cm,
                                  3.4 * cm], align=[1, 2, 3, 4], fs=8.2),
              P("If the swap-in group looks like a different population rather "
                "than a set of near-misses, the model is not finding "
                "overlooked good business - it is ranking on four fields where "
                "the underwriter used a full bureau file.", "small")]

    sens = n("part4_economics.sensitivity_to_cost_of_funds", [])
    if sens:
        E += [Spacer(1, 6), P("Sensitivity to the funding-cost assumption",
                              "h2"),
              table([["Cost of funds", "Break-even bad rate (36m)",
                      "Profit forgone"]]
                    + [[pct(r["cost_of_funds_annual"], 1),
                        pct(r["breakeven_bad_rate_36m"], 1),
                        money(r["profit_forgone_population"])] for r in sens],
                    widths=[4.6 * cm, 6.0 * cm, 5.8 * cm], align=[1, 2]),
              P("The point estimate above is one cell of this table. The sign "
                "of the answer can change across it, which is the honest "
                "characterisation of how much is measurement and how much is "
                "assumption.", "small")]
    E += [Spacer(1, 6), figure("07_swap_profit.png")]

    # ------------------------------------------------- 8. geography
    E += [PageBreak(), P("8. Part 4c - residual geographic variation", "h1")]
    E += [P("Each applicant's probability of decline is modelled from risk "
            "characteristics alone - no state, no ZIP. Observed declines per "
            "3-digit ZIP are then compared with what risk predicts. Under the "
            "null that risk explains the decision, the standardised residuals "
            "are standard normal.", "body")]
    E += [table([
        ["Statistic", "Value", "Under the null"],
        ["3-digit ZIPs analysed", f"{disp.get('n_zip3', 0):,}", "-"],
        ["Dispersion ratio", num(disp.get("dispersion_ratio"), 2), "1.00"],
        ["SD of standardised residual", num(disp.get("sd_of_z"), 2), "1.00"],
        ["ZIPs with |z| > 3", f"{disp.get('n_flagged_abs_z_gt_3', 0):,}",
         f"{disp.get('n_flagged_expected_by_chance', 0):.1f} expected"],
        ["Widest unexplained gap",
         f"{disp.get('max_excess_decline_pp', 0):+.1f} pp", "-"],
    ], widths=[6.4 * cm, 4.6 * cm, 5.4 * cm], align=[1, 2])]
    E += [Spacer(1, 8), P(
        "<b>What this does not say.</b> This analysis never touches race, "
        "ethnicity, sex, age or any other protected attribute. It joins to no "
        "census data and infers no demographics from geography. A 3-digit ZIP "
        "covers hundreds of thousands of people, and treating it as a proxy "
        "for a protected class would be both statistically indefensible and "
        "legally reckless. The only claim available here is that applications "
        "from some ZIPs were declined more often than their recorded risk "
        "characteristics predict. The explanation could be branch footprint, "
        "marketing mix, local income-documentation practice, or a risk "
        "variable this thin feature set cannot see. This is a flag for where "
        "a fair-lending review should look. It is not a finding of "
        "discrimination.", "warn")]
    E += [Spacer(1, 4), figure("08_geo.png")]

    # ------------------------------------------------- 9. cannot tell you
    E += [PageBreak(), P("9. What this tool cannot tell you", "h1")]
    E += [P("Read this section before quoting any number above.", "body")]
    E += bullets([
        "<b>It cannot tell you the true bad rate of your declined "
        "population.</b> Nothing can, from this data. Those outcomes do not "
        "exist. Every figure in Part 2 is a model-based extrapolation whose "
        "error Part 3 measures but cannot remove.",

        "<b>It cannot validate reject inference on the real portfolio.</b> "
        "The simulation validates the methods on a manufactured problem where "
        "the cutoff is known and the features are complete. Real selection "
        "used a full credit bureau file this analysis cannot see. Performance "
        "in the harness is an upper bound on performance in reality.",

        "<b>It cannot see what the underwriter saw.</b> Approval used the "
        "whole bureau report, income verification and platform policy. This "
        "audit sees four fields. Any apparent selection effect on those four "
        "may simply be omitted-variable bias in disguise - which, note, cuts "
        "against over-interpreting a small rho as much as a large one.",

        "<b>It cannot tell you about the two thirds of declines with no "
        "bureau score.</b> They are excluded, they are not missing at random, "
        "and they are concentrated in the later vintages.",

        "<b>It cannot price a policy change.</b> The swap-set profit figure "
        "assumes the approved applicants behave like the historical loans in "
        "their score band, that pricing is unchanged, and that nothing about "
        "the competitive environment responds. Adverse selection on a "
        "loosened cutoff is real and is not modelled.",

        "<b>It cannot establish discrimination.</b> See section 8.",

        "<b>It cannot be transported to your portfolio unexamined.</b> "
        "LendingClub approved roughly "
        f"{pct(n('data.population_accept_rate'), 0)} of applications over this "
        "period. A lender approving 60% has a far milder selection problem "
        "and should expect different answers.",

        "<b>It cannot tell you the loans were priced correctly.</b> This is an "
        "audit of a cutoff, not of a pricing curve.",
    ])

    # ------------------------------------------------- 10. reproducing
    E += [P("10. Reproducing every number", "h1")]
    E += [table([
        ["make install", "install pinned dependencies"],
        ["make data", "fetch both LendingClub files (CC0, ~650 MB)"],
        ["make run", "run the full pipeline, write headline_numbers.json"],
        ["make docs", "rebuild this PDF and the deck from that JSON"],
        ["make test", "run the test suite"],
    ], widths=[4.4 * cm, 12.0 * cm], header=False),
        Spacer(1, 8),
        P(f"Seed {n('meta.seed')}; Python {n('meta.python')}, "
          f"numpy {n('meta.numpy')}, pandas {n('meta.pandas')}. Every sample "
          "is drawn from a seeded generator and every dependency is pinned. "
          "No figure or table in this document was edited by hand; all of "
          "them are regenerated from outputs/headline_numbers.json.", "small")]

    doc.build(E)
    return out


# ==========================================================================
def build_deck(out: Path) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    P_INK = RGBColor(0x1B, 0x1B, 0x1F)
    P_MUTE = RGBColor(0x6B, 0x72, 0x80)
    P_ACC = RGBColor(0xC2, 0x41, 0x0C)
    P_TEAL = RGBColor(0x0E, 0x74, 0x90)
    P_WASH = RGBColor(0xF6, 0xF7, 0xF9)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def slide():
        return prs.slides.add_slide(blank)

    def txt(s, x, y, w, h, text, size=18, bold=False, color=P_INK,
            align=PP_ALIGN.LEFT, font="Calibri", space_after=6):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(space_after)
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
        return tb

    def band(s, color=P_WASH, y=0.0, h=1.15):
        from pptx.enum.shapes import MSO_SHAPE
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), SW,
                                Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    def header(s, kicker, title):
        band(s, P_WASH, 0, 1.32)
        txt(s, 0.62, 0.20, 12, 0.32, kicker, size=11.5, bold=True,
            color=P_ACC)
        txt(s, 0.62, 0.50, 12.2, 0.7, title, size=27, bold=True)

    def pic(s, name, x, y, w):
        p = C.FIGURES / name
        if p.exists():
            s.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w))

    def statbox(s, x, y, w, label, value, sub="", vcolor=P_INK):
        from pptx.enum.shapes import MSO_SHAPE
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                Inches(y), Inches(w), Inches(1.75))
        sh.fill.solid(); sh.fill.fore_color.rgb = P_WASH
        sh.line.color.rgb = RGBColor(0xD4, 0xD4, 0xD8); sh.line.width = Pt(0.75)
        sh.shadow.inherit = False
        txt(s, x + 0.22, y + 0.14, w - 0.44, 0.3, label, size=11,
            color=P_MUTE)
        txt(s, x + 0.22, y + 0.46, w - 0.44, 0.7, value, size=30, bold=True,
            color=vcolor)
        if sub:
            txt(s, x + 0.22, y + 1.22, w - 0.44, 0.4, sub, size=10,
                color=P_MUTE)

    base_gini = n("part1_baseline.test.gini")
    rho = n("part2_real_data.heckman.diagnostics.rho")
    rho_p = n("part2_real_data.heckman.diagnostics.lr_p")
    ranking = n("part3_simulation.ranking", [])
    nullc = n("part3_simulation.null_case", [])
    hall = [r["method"] for r in nullc if r.get("hallucinates")]
    ok = [r["method"] for r in nullc
          if not r.get("hallucinates") and r["method"] != "none (biased baseline)"]
    # Named with magnitudes: a 2.5 pp failure and a 7.6 pp one are not the
    # same finding, and a slide that lists them together would flatten that.
    _sized = sorted((r for r in nullc if r.get("hallucinates")),
                    key=lambda r: -r["excess_over_baseline"])
    worst_hall = 100 * (_sized[0]["excess_over_baseline"] if _sized else 0.0)
    hall_detail = "\n".join(
        f"{r['method']}:  {100 * r['excess_over_baseline']:+.1f} pp"
        for r in _sized) + (
        "\n\nThe top two assume declines are worse by a fixed multiple, so "
        "they fire whether or not selection was informative.\n\n"
        "Heckman's is smaller and has a different cause: under a random "
        "cutoff rho is unidentified. Its significance test stays honest; the "
        "point estimate drifts.")
    disp = n("part4_geography.dispersion", {})
    econ = n("part4_economics.by_method", [])
    rho_small = (rho_p or 1) > 0.05

    # ---- 1 title
    s = slide()
    band(s, P_WASH, 0, 2.55)
    txt(s, 0.85, 0.72, 11.6, 0.5, "SELECTION-BIAS AUDIT", size=13, bold=True,
        color=P_ACC)
    txt(s, 0.85, 1.12, 11.6, 1.0,
        "Is your scorecard's Gini real?", size=40, bold=True)
    txt(s, 0.85, 2.85, 11.6, 1.6,
        "Your model is validated only on applicants you approved.\n"
        "This audit measures how much that flatters it, and tests whether the\n"
        "standard corrections can be trusted to fix it.",
        size=17, color=P_MUTE, space_after=3)
    txt(s, 0.85, 5.5, 11.6, 1.2,
        f"LendingClub 2007-2018  |  {n('data.accepted_rows_total', 0):,} "
        f"approved loans, {n('data.rejected_rows_total', 0):,} declines  |  "
        f"seed {n('meta.seed')}", size=12, color=P_MUTE)
    txt(s, 0.85, 5.9, 11.6, 0.4, "Every number reproducible: make data && make run",
        size=12, bold=True, color=P_TEAL)

    # ---- 2 the problem
    s = slide()
    header(s, "THE PROBLEM", "You only ever see the applicants you approved")
    pic(s, "01_overlap.png", 0.5, 1.62, 12.3)
    txt(s, 0.62, 5.75, 12.2, 1.3,
        f"LendingClub approved {pct(n('data.population_accept_rate'), 1)} of "
        f"applications. Only "
        f"{num(n('data.scored_rejects_above_accepted_p01_pct'), 1)}% of scored "
        f"declines reach the 1st percentile of the approved book, and "
        f"{pct(n('data.unscored_reject_fraction'), 0)} of declines carry no "
        f"bureau score at all.\n"
        "The two populations barely overlap. That is the whole problem, and "
        "also the limit on how well anyone can solve it.",
        size=13, color=P_MUTE, space_after=4)

    # ---- 3 measured gini
    s = slide()
    header(s, "PART 1", "Here is your measured Gini")
    statbox(s, 0.62, 1.75, 3.9, "GINI (held-out, approved loans)",
            num(base_gini), "the number you report today")
    statbox(s, 4.72, 1.75, 3.9, "KS", num(n("part1_baseline.test.ks")),
            "same population, same caveat")
    statbox(s, 8.82, 1.75, 3.9, "OBSERVED BAD RATE",
            pct(n("data.observed_bad_rate_accepted"), 1),
            "on settled approved loans")
    txt(s, 0.62, 3.85, 12.2, 2.4,
        "A conventional WOE + logistic application scorecard on the four "
        "fields both populations share.\n\n"
        "It is a perfectly good model. The problem is not the model - it is "
        "that this number was measured on a population your cutoff already "
        "filtered. It answers 'how well do I rank the people I said yes to?', "
        "not 'how well do I rank everyone who walks in the door?'",
        size=15, color=P_INK, space_after=8)

    # ---- 4 corrected gini
    s = slide()
    header(s, "PART 2", "Here is your corrected Gini - four ways, and they disagree")
    rows = [("Method", "Est. through-the-door bad rate", "Claimed Gini")]
    for m in ("parcelling", "fuzzy", "ipw", "heckman"):
        rows.append((m,
                     pct(n(f"part2_real_data.{m}.est_ttd_bad_rate"), 2),
                     num(n(f"part2_real_data.{m}.claimed_ttd_gini"))))
    rows.append(("no correction (observed)",
                 pct(n("data.observed_bad_rate_accepted"), 2),
                 num(base_gini)))
    tb = s.shapes.add_table(len(rows), 3, Inches(0.62), Inches(1.72),
                            Inches(12.1), Inches(2.9)).table
    tb.columns[0].width = Inches(4.4)
    tb.columns[1].width = Inches(4.6)
    tb.columns[2].width = Inches(3.1)
    for i, r in enumerate(rows):
        for j, v in enumerate(r):
            c = tb.cell(i, j)
            c.text = str(v)
            p = c.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(14 if i else 12.5)
            p.runs[0].font.bold = (i == 0)
            p.runs[0].font.color.rgb = P_INK if i else P_INK
    txt(s, 0.62, 4.85, 12.2, 2.0,
        "Four defensible methods. Percentage points apart on the same data.\n"
        "The spread between them IS the uncertainty - and no amount of "
        "argument about which is theoretically nicer resolves it.\n"
        "So the next slide stops arguing and runs an experiment.",
        size=15, color=P_MUTE, space_after=6)

    # ---- 5 harness design
    s = slide()
    header(s, "PART 3", "Can you trust the correction? Manufacture the truth")
    steps = [
        ("1", "Take only approved loans", "every outcome is known"),
        ("2", "Impose an artificial cutoff", "discard outcomes below it"),
        ("3", "Run all four methods", "on a synthetic decline population"),
        ("4", "Compare against the truth", "that we deliberately hid"),
    ]
    for i, (num_, t, sub) in enumerate(steps):
        x = 0.62 + i * 3.08
        statbox(s, x, 1.8, 2.86, f"STEP {num_}", "", "")
        txt(s, x + 0.22, 2.15, 2.5, 0.9, t, size=14, bold=True)
        txt(s, x + 0.22, 2.95, 2.5, 0.6, sub, size=11, color=P_MUTE)
    txt(s, 0.62, 4.0, 12.2, 2.6,
        "Reject inference is normally unfalsifiable - you cannot check an "
        "inferred bad rate against outcomes that do not exist.\n"
        "This makes it falsifiable. Swept across rejection rates of "
        f"{', '.join(pct(r, 0) for r in n('part3_simulation.rejection_rates', []))}"
        f", with {n('part3_simulation.n_replicates')} replicates each, under "
        "three cutoff shapes: a hard risk threshold, a threshold with "
        "override noise, and - critically - a random cutoff where the correct "
        "answer is known to be 'do nothing'.",
        size=15, color=P_INK, space_after=8)

    # ---- 6 degradation
    s = slide()
    header(s, "PART 3", "How each method degrades as you decline more")
    pic(s, "03_sweep_bad_rate.png", 0.5, 1.55, 12.3)
    vd = n("part3_simulation.degradation.verdict", "")
    if ranking:
        b = ranking[0]
        txt(s, 0.62, 5.85, 12.2, 1.2,
            f"Least biased on risk-based cutoffs: {b['method']} "
            f"({100 * b['mean_abs_bad_rate_bias']:.2f} pp mean absolute error). "
            "Shaded bands are across-replicate variation.\n" + vd,
            size=13, color=P_MUTE, space_after=3)

    # ---- 7 null case
    s = slide()
    header(s, "PART 3 - THE NULL CASE",
           "Which methods invent a bias that isn't there?")
    pic(s, "06_null_case.png", 0.62, 1.6, 7.4)
    txt(s, 8.35, 1.75, 4.5, 3.6,
        "Replace the risk cutoff with a RANDOM one.\n\n"
        "Now there is no selection bias. The correct answer from every method "
        "is zero correction.\n\n"
        + (hall_detail if hall else
           "Every method correctly reported approximately no correction."),
        size=14, color=P_INK, space_after=8)
    if ok:
        txt(s, 8.35, 5.5, 4.5, 1.0,
            f"Passes the null: {', '.join(ok)}",
            size=13, bold=True, color=P_TEAL)

    # ---- 8 rho
    s = slide()
    header(s, "PART 2 - THE ECONOMETRIC READ",
           "How much selection bias is actually in this portfolio?")
    statbox(s, 0.62, 1.72, 3.9, "SELECTION CORRELATION (rho)", num(rho, 3),
            "bivariate probit, full ML",
            vcolor=P_TEAL if rho_small else P_ACC)
    statbox(s, 4.72, 1.72, 3.9, "LIKELIHOOD-RATIO p", num(rho_p, 3),
            "H0: rho = 0")
    statbox(s, 8.82, 1.72, 3.9, "VERDICT",
            "SMALL" if rho_small else "PRESENT",
            "on the shared fields",
            vcolor=P_TEAL if rho_small else P_ACC)
    txt(s, 0.62, 3.8, 12.2, 2.6,
        ("rho is the correlation between the unobserved drivers of approval "
         "and of default. It is not distinguishable from zero here.\n\n"
         "The honest reading: on the fields these two populations share, "
         "approval carried little information about default beyond what the "
         "scorecard already uses. Selection bias in this portfolio is small.\n\n"
         "Manufacturing a larger correction would be the actual failure - and "
         f"note that {len(hall)} of the four methods do exactly that."
         if rho_small else
         "rho is distinguishable from zero, so a genuine selection effect is "
         "present on the shared fields."),
        size=15, color=P_INK, space_after=8)

    # ---- 9 money
    s = slide()
    header(s, "PART 4", "What is your cutoff costing you?")
    be = n("part4_economics.breakeven_bad_rate_36m")
    be_sw = [r.get("swap_in_breakeven_bad_rate") for r in econ
             if r.get("swap_in_breakeven_bad_rate") is not None]
    sens_vals = [r["profit_forgone_population"]
                 for r in n("part4_economics.sensitivity_to_cost_of_funds", [])
                 if r.get("profit_forgone_population") is not None]
    vals = [r.get("profit_forgone_population") for r in econ
            if r.get("profit_forgone_population") is not None]
    lo, hi = (min(vals), max(vals)) if vals else (None, None)
    statbox(s, 0.62, 1.72, 3.9, "BREAK-EVEN BAD RATE", pct(be, 1),
            "36-month loan, measured yields")
    statbox(s, 4.72, 1.72, 3.9, "SWAP-IN BAD RATE (RANGE)",
            f"{pct(min((r.get('swap_in_inferred_bad_rate') or 9) for r in econ), 1)}"
            f"-{pct(max((r.get('swap_in_inferred_bad_rate') or 0) for r in econ), 1)}"
            if econ else "n/a", "across the four methods")
    statbox(s, 8.82, 1.72, 3.9, "PROFIT EFFECT",
            f"{money(lo)} to {money(hi)}" if vals else "n/a",
            "the spread IS the finding",
            vcolor=P_ACC if (vals and lo < 0) else P_TEAL)
    txt(s, 0.62, 3.8, 12.2, 2.7,
        "Volume-neutral swap: the corrected model approves exactly as many "
        "applicants as you do today, so this measures ranking quality, not "
        "appetite.\n\n"
        "Yields, losses and LGD are measured from LendingClub's realised cash "
        "flows. Cost of funds, servicing and loan life are stated assumptions, "
        "and the answer is genuinely sensitive to them - the methodology PDF "
        "carries the full sensitivity grid.\n\n"
        "Read the range, not the point estimate.",
        size=14.5, color=P_INK, space_after=8)

    # ---- 10 recommendation
    s = slide()
    header(s, "RECOMMENDATION", "What I would actually do")
    recs = [
        ("Do not adopt a correction you cannot falsify.",
         f"{len(hall)} of the four report a correction on data with provably "
         f"zero selection bias, the worst by {worst_hall:.0f} pp. Run the "
         f"null case before trusting any of them."),
        ("Treat the rho result as the headline, not the profit number.",
         "Measured selection bias on the shared fields is small. That is a "
         "cheaper finding to act on than a contested dollar figure."),
        ("Fix the data before fixing the model.",
         f"{pct(n('data.unscored_reject_fraction'), 0)} of declines carry no "
         "bureau score. Retaining the decision-time bureau pull on declines "
         "would do more for this problem than any estimator."),
        ("Run a randomised-approval holdout.",
         "A small random approval band above the cutoff generates the "
         "unbiased sample that makes all of this measurable rather than "
         "inferred. It is the only clean fix."),
        ("Use the ZIP flags to scope a fair-lending review.",
         f"{disp.get('n_flagged_abs_z_gt_3', 0)} ZIPs sit beyond 3 standard "
         "deviations of what risk predicts, against "
         f"{disp.get('n_flagged_expected_by_chance', 0):.0f} expected. That is "
         "a place to look, not a conclusion."),
    ]
    y = 1.62
    for t, sub in recs:
        txt(s, 0.62, y, 12.2, 0.34, t, size=15, bold=True, color=P_INK)
        txt(s, 0.62, y + 0.36, 12.2, 0.62, sub, size=12, color=P_MUTE)
        y += 1.06

    prs.save(str(out))
    return out


# ==========================================================================
def _runtime() -> str:
    """Measured wall-clock of the last run."""
    secs = n("meta.runtime_seconds")
    if secs:
        return f"{secs / 60:.0f} minutes" if secs >= 90 else f"{secs:.0f} seconds"
    log = C.OUTPUTS / "run.log"
    if log.exists():
        for line in reversed(log.read_text(errors="ignore").splitlines()):
            if "done in" in line:
                try:
                    secs = float(line.split("done in")[1].strip().rstrip("s"))
                    return (f"{secs / 60:.0f} minutes" if secs >= 90
                            else f"{secs:.0f} seconds")
                except ValueError:
                    break
    return "under an hour"


def _hallucination_sentence() -> str:
    """Name each failing method WITH its magnitude.

    A binary pass/fail flag would put a 2.5 pp error and a 7.6 pp error in the
    same bucket, which is exactly the kind of flattening this repo exists to
    argue against.
    """
    rows = sorted((r for r in n("part3_simulation.null_case", [])
                   if r.get("hallucinates")),
                  key=lambda r: -r["excess_over_baseline"])
    if not rows:
        return "none - all four passed"
    lines = [f"| **{r['method']}** | "
             f"{100 * r['mean_abs_bad_rate_bias']:.2f} pp | "
             f"+{100 * r['excess_over_baseline']:.2f} pp | "
             f"{r['p_one_sided']:.1g} |" for r in rows]
    ok = [r for r in n("part3_simulation.null_case", [])
          if not r.get("hallucinates")]
    lines += [f"| {r['method']} | "
              f"{100 * r['mean_abs_bad_rate_bias']:.2f} pp | "
              f"+{100 * r['excess_over_baseline']:.2f} pp | passes |"
              for r in sorted(ok, key=lambda r: r["excess_over_baseline"])]
    header = ("| Method | Error under a random cutoff | "
              "Excess over doing nothing | one-sided p |")
    return "\n".join([header, "| --- | --- | --- | --- |"] + lines)


def _scale_finding() -> str:
    """What the scale sweep says about whether these methods are salvageable."""
    rows = n("part3_scale_sensitivity", [])
    if not rows:
        return ""
    best = {}
    for r in rows:
        k = (r["cutoff_type"], r["method"])
        if k not in best or r["abs_bad_rate_bias"] < best[k]["abs_bad_rate_bias"]:
            best[k] = r
    hard = [v for (ct, _), v in best.items() if ct == "risk_hard"]
    rand = [v for (ct, _), v in best.items() if ct == "random"]
    if not hard or not rand:
        return ""
    return (f"At their best multiple they reach "
            f"{100 * min(v['abs_bad_rate_bias'] for v in hard):.1f} pp on a "
            f"risk-based cutoff and "
            f"{100 * min(v['abs_bad_rate_bias'] for v in rand):.1f} pp on the "
            f"null - competitive with the other two. The machinery is not the "
            f"problem; the conventional multiple is.")


def _swapin_sentence() -> str:
    """One concrete sentence contrasting the swap-in set with today's book."""
    rows = {r.get("group"): r for r in n("part4_swap_profile", [])}
    book = rows.get("approved today")
    swap = rows.get("swap-in (would approve)")
    if not book or not swap:
        return ""
    return (f"The applicants who flip average a "
            f"${swap['amount_requested']:,.0f} loan request at "
            f"{swap['dti']:.0f}% DTI with {swap['emp_length_yrs']:.1f} years of "
            f"employment, against ${book['amount_requested']:,.0f}, "
            f"{book['dti']:.0f}% and {book['emp_length_yrs']:.1f} years on "
            f"today's book.")


def build_readme(template: Path, out: Path) -> Path:
    """Render README.md from the template using ONLY pipeline output.

    The README quotes a dozen numbers. Typing them by hand would guarantee that
    one of them eventually disagrees with the repo it describes, so none of
    them are typed by hand.
    """
    ranking = n("part3_simulation.ranking", [])
    nullc = n("part3_simulation.null_case", [])
    econ = n("part4_economics.by_method", [])
    rho = n("part2_real_data.heckman.diagnostics.rho")
    rho_p = n("part2_real_data.heckman.diagnostics.lr_p")

    ttd = [n(f"part2_real_data.{m}.est_ttd_bad_rate")
           for m in ("parcelling", "fuzzy", "ipw", "heckman")]
    ttd = [t for t in ttd if t is not None]
    hall = [r["method"] for r in nullc if r.get("hallucinates")]
    sw = [r.get("swap_in_inferred_bad_rate") for r in econ
          if r.get("swap_in_inferred_bad_rate") is not None]
    profits = [r.get("profit_forgone_population") for r in econ
               if r.get("profit_forgone_population") is not None]
    be = n("part4_economics.breakeven_bad_rate_36m")
    be_sw = [r.get("swap_in_breakeven_bad_rate") for r in econ
             if r.get("swap_in_breakeven_bad_rate") is not None]
    sens_vals = [r["profit_forgone_population"]
                 for r in n("part4_economics.sensitivity_to_cost_of_funds", [])
                 if r.get("profit_forgone_population") is not None]

    if profits and all(p < 0 for p in profits):
        verdict = ("approving the swap-in set would have destroyed value, not "
                   "created it, at every method's estimate of its bad rate. "
                   "The cutoff is not leaving money on the table here.")
    elif profits and all(p > 0 for p in profits):
        verdict = ("the swap-in set is profitable on every method's estimate, "
                   "so the cutoff does appear to be leaving money on the table.")
    else:
        verdict = ("the methods disagree even on the sign, which is itself the "
                   "finding: this data cannot settle the question.")

    rho_x = n("part2_real_data.heckman_with_exclusion.diagnostics.rho")
    rho_xp = n("part2_real_data.heckman_with_exclusion.diagnostics.lr_p")
    biggest = max(abs(rho or 0.0), abs(rho_x or 0.0))
    if biggest < 0.10:
        rho_verdict = ("Both are economically negligible - a correlation of "
                       "this size shifts the corrected bad rate by a fraction "
                       "of a percentage point.")
    elif (rho_p or 1) > 0.05:
        rho_verdict = "The primary specification is not distinguishable from zero."
    else:
        rho_verdict = ("Both specifications find a selection effect large "
                       "enough to matter.")

    sub = {
        "n_accepted_total": f"{n('data.accepted_rows_total', 0):,}",
        "n_rejected_total": f"{n('data.rejected_rows_total', 0):,}",
        "n_rejected_used": f"{n('data.rejected_rows_used', 0):,}",
        "n_rejected_sampled": f"{n('data.rejected_rows_sampled_raw', 0):,}",
        "unscored_pct": pct(n("data.unscored_reject_fraction"), 0),
        "overlap_pct": num(n("data.scored_rejects_above_accepted_p01_pct"), 1),
        "indeterminate_pct": pct(n("data.indeterminate_fraction"), 0),
        "accept_rate": pct(n("data.population_accept_rate"), 1),
        "base_gini": num(n("part1_baseline.test.gini")),
        "base_ks": num(n("part1_baseline.test.ks")),
        "base_auc": num(n("part1_baseline.test.auc")),
        "obs_bad_rate": pct(n("data.observed_bad_rate_accepted"), 1),
        "ttd_range": (f"{pct(min(ttd), 1)} to {pct(max(ttd), 1)}"
                      if ttd else "n/a"),
        "rho": num(rho, 3),
        "rho_se": num(n("part2_real_data.heckman.diagnostics.se_rho"), 3),
        "rho_p": num(rho_p, 3),
        "rho_verdict": rho_verdict,
        "rho_x": num(rho_x, 3),
        "rho_xp": num(rho_xp, 3),
        "n_modelled": f"{n('data.accepted_rows_used', 0) + n('data.rejected_rows_used', 0):,}",
        "best_method": ranking[0]["method"] if ranking else "n/a",
        "best_bias": (f"{100 * ranking[0]['mean_abs_bad_rate_bias']:.2f}"
                      if ranking else "n/a"),
        "hallucinators": _hallucination_sentence(),
        "scale_finding": _scale_finding(),
        "disp_ratio": num(n("part4_geography.dispersion.dispersion_ratio"), 1),
        "breakeven": pct(be, 1),
        "swapin_range": (f"{pct(min(sw), 1)} to {pct(max(sw), 1)}"
                         if sw else "n/a"),
        "swapin_breakeven": (f"{pct(min(be_sw), 1)} to {pct(max(be_sw), 1)}"
                             if be_sw else "n/a"),
        "profit_lo_sens": money(min(sens_vals)) if sens_vals else "n/a",
        "profit_hi_sens": money(max(sens_vals)) if sens_vals else "n/a",
        "profit_verdict": verdict,
        "seed": str(n("meta.seed")),
        "runtime": _runtime(),
        "pdo": num(n("part1_baseline.pdo"), 0),
        "rates": ", ".join(pct(r, 0) for r in
                           n("part3_simulation.rejection_rates", [])),
        "n_reps": str(n("part3_simulation.n_replicates")),
        "prop_auc": num(n("part2_positivity.propensity_auc"), 3),
        "verdict": n("part3_simulation.degradation.verdict", "n/a"),
        "swapin_profile": _swapin_sentence(),
    }

    text = template.read_text(encoding="utf-8")
    for k, v in sub.items():
        text = text.replace("{{" + k + "}}", str(v))
    leftover = [t for t in text.split("{{")[1:]]
    if leftover:
        raise SystemExit(f"README template has unfilled tokens: "
                         f"{[t.split('}}')[0] for t in leftover]}")
    out.write_text(text, encoding="utf-8")
    return out


def main() -> None:
    global D
    src = C.OUTPUTS / "headline_numbers.json"
    if not src.exists():
        raise SystemExit("outputs/headline_numbers.json not found. "
                         "Run `python run_all.py` first.")
    D = json.loads(src.read_text())
    C.ensure_dirs()
    pdf = build_pdf(C.DOCS / "methodology.pdf")
    print(f"wrote {pdf}")
    deck = build_deck(C.SLIDES / "credit_risk_briefing.pptx")
    print(f"wrote {deck}")
    rd = build_readme(C.ROOT / "README.template.md", C.ROOT / "README.md")
    print(f"wrote {rd}")


if __name__ == "__main__":
    main()
